import asyncio
import hashlib
import json
import os
import time
import sqlite3
from typing import List, Optional

import numpy as np
from PIL import Image, ExifTags
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

from .config import load_config
from .db import get_db
from .embeddings import get_image_embedding, get_sbert_embedding
from .animal_detector import detect_objects
from .ocr import run_ocr
from .search_index import build_index, load_index, save_index

SUPPORTED_IMAGE_TYPES = {".jpg", ".jpeg", ".png", ".heic", ".webp", ".tiff", ".tif", ".bmp", ".gif", ".ico", ".svg"}
SUPPORTED_VIDEO_TYPES = {".mp4", ".mov", ".avi", ".mkv", ".flv", ".wmv", ".webm", ".m4v", ".3gp", ".ts"}
SUPPORTED_AUDIO_TYPES = {".mp3", ".wav", ".flac", ".aac", ".m4a", ".ogg", ".wma", ".opus", ".alac", ".aif", ".aiff", ".ape", ".dsd"}
SUPPORTED_TEXT_TYPES = {".pdf", ".docx", ".doc", ".txt", ".md", ".rtf", ".odt", ".xls", ".xlsx", ".csv", ".pptx", ".ppt", ".json", ".xml", ".html", ".htm"}

# File type categories for filtering
FILE_TYPE_CATEGORIES = {
    "images": SUPPORTED_IMAGE_TYPES,
    "videos": SUPPORTED_VIDEO_TYPES,
    "audio": SUPPORTED_AUDIO_TYPES,
    "text": SUPPORTED_TEXT_TYPES,
}


def convert_aif_to_wav(aif_path: str) -> str:
    """Convert AIF file to WAV using ffmpeg (memory efficient, streaming). Returns path to WAV file."""
    import subprocess
    import hashlib
    
    try:
        ext = os.path.splitext(aif_path)[1].lower()
        if ext not in ['.aif', '.aiff']:
            return aif_path  # Not an AIF file
        
        # Create cache directory
        cache_dir = os.path.join(os.path.dirname(aif_path), '.audio-cache')
        os.makedirs(cache_dir, exist_ok=True)
        
        # Generate cache filename
        file_hash = hashlib.md5(aif_path.encode()).hexdigest()
        wav_path = os.path.join(cache_dir, f'{file_hash}.wav')
        
        # Check if already converted
        if os.path.exists(wav_path) and os.path.getsize(wav_path) > 0:
            print(f"[AUDIO] ✓ Using cached conversion: {os.path.basename(wav_path)}")
            return wav_path
        
        print(f"[AUDIO] Converting AIF to WAV: {os.path.basename(aif_path)}")
        
        # Use ffmpeg for efficient streaming conversion (no memory load)
        try:
            result = subprocess.run(
                ['ffmpeg', '-i', aif_path, '-acodec', 'pcm_s16le', '-ar', '44100',
                 wav_path, '-y'],
                capture_output=True,
                timeout=300,  # 5 minute timeout for large files
                check=False
            )
            
            if result.returncode == 0 and os.path.exists(wav_path):
                file_size = os.path.getsize(wav_path)
                if file_size > 0:
                    print(f"[AUDIO] ✓ Conversion successful: {file_size} bytes")
                    return wav_path
                else:
                    print(f"[AUDIO] ✗ Empty WAV file created")
                    return aif_path
            else:
                print(f"[AUDIO] ✗ ffmpeg conversion failed: {result.stderr.decode()}")
                return aif_path
                
        except subprocess.TimeoutExpired:
            print(f"[AUDIO] ✗ Conversion timeout (file too large or slow disk)")
            return aif_path
        except FileNotFoundError:
            print(f"[AUDIO] ✗ ffmpeg not found - cannot convert AIF files")
            return aif_path
        
    except Exception as e:
        print(f"[AUDIO] ✗ Conversion error: {e}")
        return aif_path


def md5sum(path: str, chunk_size: int = 1024 * 1024) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        while True:
            chunk = f.read(chunk_size)
            if not chunk:
                break
            h.update(chunk)
    return h.hexdigest()


def is_media(path: str, skip_videos: bool) -> bool:
    ext = os.path.splitext(path.lower())[1]
    if ext in SUPPORTED_IMAGE_TYPES:
        return True
    if not skip_videos and ext in SUPPORTED_VIDEO_TYPES:
        return True
    if ext in SUPPORTED_AUDIO_TYPES:
        return True
    if ext in SUPPORTED_TEXT_TYPES:
        return True
    return False


def extract_text_content(path: str) -> str:
    """Extract text content from various document formats."""
    ext = os.path.splitext(path.lower())[1]
    text = ""
    
    try:
        # Plain text files
        if ext in ['.txt', '.md']:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        
        # JSON files
        elif ext == '.json':
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    data = json.load(f)
                    text = json.dumps(data, indent=2)
            except:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
        
        # XML and HTML files
        elif ext in ['.xml', '.html', '.htm']:
            try:
                from html.parser import HTMLParser
                
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text_parts = []
                    
                    def handle_data(self, data):
                        text_data = data.strip()
                        if text_data:
                            self.text_parts.append(text_data)
                
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    html_content = f.read()
                
                parser = TextExtractor()
                parser.feed(html_content)
                text = ' '.join(parser.text_parts)
            except:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
        
        # PDF files
        elif ext == '.pdf':
            try:
                import PyPDF2
                with open(path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    text_parts = []
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                    text = ' '.join(text_parts)
            except ImportError:
                print(f"[TEXT] Warning: PyPDF2 not available for {path}")
            except Exception as e:
                print(f"[TEXT] Error reading PDF {path}: {e}")
        
        # Word documents (.docx)
        elif ext == '.docx':
            try:
                from docx import Document
                doc = Document(path)
                text_parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_parts.append(cell.text)
                text = ' '.join(text_parts)
            except ImportError:
                print(f"[TEXT] Warning: python-docx not available for {path}")
            except Exception as e:
                print(f"[TEXT] Error reading DOCX {path}: {e}")
        
        # Old Word documents (.doc)
        elif ext == '.doc':
            try:
                import docx2txt
                text = docx2txt.process(path)
            except ImportError:
                print(f"[TEXT] Warning: docx2txt not available for {path}")
            except Exception as e:
                print(f"[TEXT] Error reading DOC {path}: {e}")
        
        # Excel files (.xlsx)
        elif ext == '.xlsx':
            try:
                import openpyxl
                wb = openpyxl.load_workbook(path)
                text_parts = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.value:
                                text_parts.append(str(cell.value))
                text = ' '.join(text_parts)
            except ImportError:
                print(f"[TEXT] Warning: openpyxl not available for {path}")
            except Exception as e:
                print(f"[TEXT] Error reading XLSX {path}: {e}")
        
        # Old Excel files (.xls)
        elif ext == '.xls':
            try:
                import xlrd
                workbook = xlrd.open_workbook(path)
                text_parts = []
                for sheet in workbook.sheets():
                    for row in range(sheet.nrows):
                        for col in range(sheet.ncols):
                            cell_value = sheet.cell(row, col).value
                            if cell_value:
                                text_parts.append(str(cell_value))
                text = ' '.join(text_parts)
            except ImportError:
                print(f"[TEXT] Warning: xlrd not available for {path}")
            except Exception as e:
                print(f"[TEXT] Error reading XLS {path}: {e}")
        
        # CSV files
        elif ext == '.csv':
            try:
                import csv
                text_parts = []
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    csv_reader = csv.reader(f)
                    for row in csv_reader:
                        text_parts.extend([cell for cell in row if cell])
                text = ' '.join(text_parts)
            except Exception as e:
                print(f"[TEXT] Error reading CSV {path}: {e}")
        
        # PowerPoint files (.pptx)
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            text_parts.append(shape.text)
                text = ' '.join(text_parts)
            except ImportError:
                print(f"[TEXT] Warning: python-pptx not available for {path}")
            except Exception as e:
                print(f"[TEXT] Error reading PPTX {path}: {e}")
        
        # Old PowerPoint files (.ppt)
        elif ext == '.ppt':
            try:
                from pptx import Presentation
                # Try to open as PPTX first (Office 2007+)
                prs = Presentation(path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            text_parts.append(shape.text)
                text = ' '.join(text_parts)
            except ImportError:
                print(f"[TEXT] Warning: python-pptx not available for {path}")
            except Exception as e:
                print(f"[TEXT] Error reading PPT {path}: {e}")
        
        # RTF files
        elif ext == '.rtf':
            try:
                from striprtf.striprtf import rtf_to_text
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    rtf_content = f.read()
                    text = rtf_to_text(rtf_content)
            except ImportError:
                # Fallback: just read as text
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
            except Exception as e:
                print(f"[TEXT] Error reading RTF {path}: {e}")
        
        # ODF files (.odt)
        elif ext == '.odt':
            try:
                from odf import opendocument, text as odf_text
                doc = opendocument.load(path)
                text_parts = []
                for paragraph in doc.getElementsByType(odf_text.P):
                    for node in paragraph.childNodes:
                        if node.nodeType == node.TEXT_NODE:
                            text_parts.append(str(node.data))
                text = ' '.join(text_parts)
            except ImportError:
                print(f"[TEXT] Warning: odfpy not available for {path}")
            except Exception as e:
                print(f"[TEXT] Error reading ODT {path}: {e}")
    
    except Exception as e:
        print(f"[TEXT] Unexpected error extracting text from {path}: {e}")
    
    # Return cleaned text (limit to first 50000 chars to avoid memory issues)
    return text.strip()[:50000] if text else ""


def extract_exif(path: str) -> dict:
    meta = {
        "date_taken": None,
        "camera": None,
        "lens": None,
        "width": None,
        "height": None,
    }
    try:
        with Image.open(path) as img:
            meta["width"], meta["height"] = img.size
            info = img._getexif() or {}
            tag_map = {ExifTags.TAGS.get(k, k): v for k, v in info.items()}
            if "DateTimeOriginal" in tag_map:
                dt = tag_map["DateTimeOriginal"]
                meta["date_taken"] = int(time.mktime(time.strptime(dt, "%Y:%m:%d %H:%M:%S")))
            meta["camera"] = tag_map.get("Model")
            meta["lens"] = tag_map.get("LensModel")
    except Exception:
        pass
    return meta


def to_blob(vec: List[float]) -> sqlite3.Binary:
    if vec is None:
        return None
    return sqlite3.Binary(np.asarray(vec, dtype="float32").tobytes())


def from_blob(blob) -> List[float]:
    """Deserialize embedding from binary blob."""
    if blob is None:
        return None
    return np.frombuffer(blob, dtype="float32").tolist()


def store_object_detections(conn: sqlite3.Connection, media_id: int, detections: list):
    now = int(time.time())
    conn.execute("DELETE FROM object_detections WHERE media_id = ?", (media_id,))
    for det in detections:
        conn.execute(
            """
            INSERT INTO object_detections (media_id, class_name, confidence, bbox, class_id, source, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                media_id,
                det["class_name"],
                det.get("confidence"),
                json.dumps(det.get("bbox")),
                det.get("class_id"),
                det.get("source", "yolo"),
                now,
                now,
            ),
        )


def store_face_embeddings(conn: sqlite3.Connection, media_id: int, faces: list):
    now = int(time.time())
    print(f"[INDEXING] store_face_embeddings called with {len(faces)} faces for media_id={media_id}")
    conn.execute("DELETE FROM face_embeddings WHERE media_id = ?", (media_id,))
    if not faces:
        # IMPORTANT: Don't insert marker rows during initial indexing!
        # The batch face detection phase checks for missing entries to know which files need processing.
        # If we insert markers here, batch face detection will skip these files.
        print(f"[INDEXING]   No faces to store for media_id={media_id} - skipping (batch face detection will process later)")
        return
    else:
        for i, face in enumerate(faces):
            embedding = face.get("embedding")
            # DEBUG: Log the embedding info
            if embedding:
                print(f"[INDEXING]   Face {i}: Storing embedding with {len(embedding)} floats, score={face.get('score')}")
            else:
                print(f"[INDEXING]   Face {i}: WARNING - No embedding! Keys: {list(face.keys())}, Values: {face}")
            conn.execute(
                """
                INSERT INTO face_embeddings (media_id, embedding, bbox, confidence, label, created_at, updated_at)
                VALUES (?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    media_id,
                    to_blob(embedding),
                    json.dumps(face.get("bbox")),
                    face.get("score"),
                    face.get("label"),
                    now,
                    now,
                ),
            )


def store_ocr_results(conn: sqlite3.Connection, media_id: int, results: list):
    now = int(time.time())
    conn.execute("DELETE FROM ocr_results WHERE media_id = ?", (media_id,))
    for res in results:
        conn.execute(
            """
            INSERT INTO ocr_results (media_id, text, confidence, bbox, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?)
            """,
            (
                media_id,
                res.get("text"),
                res.get("confidence"),
                json.dumps(res.get("bbox")),
                now,
                now,
            ),
        )


def store_uncertain_detections(
    conn: sqlite3.Connection,
    media_id: int,
    detection_type: str,
    detections: list,
) -> int:
    """Store uncertain detections for later user review."""
    count = 0
    now = int(time.time())
    for det in detections:
        conn.execute(
            """
            INSERT INTO uncertain_detections 
            (media_id, detection_type, class_name, confidence, bbox, raw_data, reviewed, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                media_id,
                detection_type,
                det.class_name if hasattr(det, "class_name") else None,
                det.confidence if hasattr(det, "confidence") else None,
                json.dumps(det.bbox) if hasattr(det, "bbox") else None,
                json.dumps({
                    "class_name": det.class_name if hasattr(det, "class_name") else None,
                    "confidence": det.confidence if hasattr(det, "confidence") else None,
                    "class_id": det.class_id if hasattr(det, "class_id") else None,
                }),
                0,
                now,
                now,
            ),
        )
        count += 1
    return count


def upsert_media(conn: sqlite3.Connection, record: dict) -> int:
    path = record.get('path')
    print(f"[UPSERT] Attempting to insert media: {path}", flush=True)
    try:
        conn.execute(
            """
            INSERT INTO media_files (path, hash, type, date_taken, location, size, width, height, camera, lens, text_embedding, clip_embedding, objects, faces, animals, text_content, created_at, updated_at)
            VALUES (:path, :hash, :type, :date_taken, :location, :size, :width, :height, :camera, :lens, :text_embedding, :clip_embedding, :objects, :faces, :animals, :text_content, :created_at, :updated_at)
            ON CONFLICT(path) DO UPDATE SET
                hash=excluded.hash,
                type=excluded.type,
                date_taken=excluded.date_taken,
                location=excluded.location,
                size=excluded.size,
                width=excluded.width,
                height=excluded.height,
                camera=excluded.camera,
                lens=excluded.lens,
                text_embedding=excluded.text_embedding,
                clip_embedding=excluded.clip_embedding,
                objects=excluded.objects,
                faces=excluded.faces,
                animals=excluded.animals,
                text_content=excluded.text_content,
                updated_at=excluded.updated_at;
            """,
            record,
        )
        cur = conn.execute("SELECT id FROM media_files WHERE path=?", (path,))
        row = cur.fetchone()
        media_id = row[0] if row else -1
        print(f"[UPSERT_SUCCESS] Media ID: {media_id} for {path}", flush=True)
        return media_id
    except Exception as e:
        print(f"[UPSERT_ERROR] Failed to upsert {path}: {e}", flush=True)
        import traceback
        traceback.print_exc()
        raise


async def index_all_sources(media_paths: list, skip_videos: bool = False) -> int:
    """Index all configured media sources, maintaining aggregate progress across all folders."""
    from .main import _get_silo_indexing_state, _currently_processing_silo
    from .silo_manager import SiloManager
    from .db import get_db_path
    
    indexing_state = _get_silo_indexing_state()
    count = 0
    cfg = load_config()
    import gc
    
    # Show which database we're using
    active_silo = SiloManager.get_active_silo()
    silo_name = active_silo.get("name", "unknown") if active_silo else "unknown"
    db_path = get_db_path()
    
    print(f"[INDEXING] ==========================================")
    print(f"[INDEXING] UNIFIED INDEXING: All sources into silo: {silo_name}")
    print(f"[INDEXING] Database: {db_path}")
    print(f"[INDEXING] ==========================================")
    
    # Collect ALL files from all paths upfront
    print(f"[INDEXING] Scanning all sources for media files...")
    all_files = []
    files_by_source = {}
    
    for media_path in media_paths:
        if not os.path.exists(media_path):
            print(f"[INDEXING] ⚠ Source does not exist: {media_path}")
            continue
        
        source_files = []
        for dirpath, _, filenames in os.walk(media_path):
            for name in filenames:
                full = os.path.join(dirpath, name)
                if is_media(full, skip_videos):
                    all_files.append(full)
                    source_files.append(full)
        
        files_by_source[media_path] = source_files
        print(f"[INDEXING] Source '{media_path}': {len(source_files)} media files")
    
    total_files = len(all_files)
    print(f"[INDEXING] Total across all sources: {total_files} files")
    
    if not all_files:
        print(f"[INDEXING] No media files found in any source!")
        indexing_state["status"] = "complete"
        indexing_state["percentage"] = 100
        indexing_state["current_file"] = "No media files found to index"
        return 0
    
    # Update indexing state with total
    indexing_state["total"] = total_files
    indexing_state["total_files_found"] = total_files
    
    clip_embeddings = []
    clip_ids = []
    total_faces = 0
    total_animals = 0
    already_indexed_count = 0
    needs_indexing_count = 0
    
    # Process ALL files in order
    for idx, full in enumerate(all_files):
        # Check if processing is paused
        from .main import _processing_paused
        if _processing_paused:
            print(f"[INDEXING] Processing is paused. Waiting...")
            await asyncio.sleep(2)
            idx -= 1  # Retry this file
            continue
        
        name = os.path.basename(full)
        file_hash = md5sum(full)
        
        # Check if already indexed (brief DB query, close immediately)
        with get_db() as conn:
            cur = conn.execute("SELECT id FROM media_files WHERE hash = ?", (file_hash,))
            result = cur.fetchone()
            if result:
                # File already indexed - skip it
                already_indexed_count += 1
                indexing_state["already_indexed"] = already_indexed_count
                indexing_state["processed"] = indexing_state.get("processed", 0) + 1
                # Also increment aggregate counter
                indexing_state["processed_count"] = indexing_state.get("processed_count", 0) + 1
                
                # Update aggregate progress
                processed = indexing_state.get("processed_count", 0)
                percentage = int((processed / max(1, total_files)) * 100)
                remaining = max(0, total_files - processed)
                
                # Update percentage in state for logs
                indexing_state["percentage"] = percentage
                
                # Show breakdown by type
                by_type = indexing_state.get("by_type", {})
                type_info = f"Images: {by_type.get('images', 0)} | Videos: {by_type.get('videos', 0)} | Audio: {by_type.get('audio', 0)} | Docs: {by_type.get('text', 0)}"
                
                indexing_state["current_file"] = (
                    f"📊 {processed}/{total_files} ({percentage}%) | {type_info} | "
                    f"Remaining: {remaining} | Skipped: {name}"
                )
                
                print(f"[INDEXING] [{idx+1}/{total_files}] ⊘ SKIPPING (already in DB): {name} ({processed}/{total_files}, {percentage}%)")
                await asyncio.sleep(0.1)
                continue
            else:
                print(f"[INDEXING] [{idx+1}/{total_files}] Hash NOT found in DB: {file_hash[:16]}... for {name}")
        
        # File needs indexing - mark it
        needs_indexing_count += 1
        indexing_state["needs_indexing"] = needs_indexing_count
        
        print(f"[INDEXING] [{idx+1}/{total_files}] NEW FILE (will index): {name}")
        print(f"[INDEXING]   Stats: {already_indexed_count} indexed | {needs_indexing_count} new")
        
        try:
            stat = os.stat(full)
            meta = extract_exif(full)
            ext = os.path.splitext(full)[1].lower()

            print(f"[INDEXING]   CLIP embedding...")
            clip_embed = get_image_embedding(full) if ext in SUPPORTED_IMAGE_TYPES else None
            await asyncio.sleep(0.05)

            # Object detection will run AFTER indexing completes in batch mode via worker
            # Skip per-file object detection during indexing to prevent crashes
            objects = []
            animals = []
            objects_json = json.dumps([
                {
                    "class": o.class_name,
                    "confidence": o.confidence,
                    "bbox": o.bbox,
                    "class_id": o.class_id,
                    "is_animal": o.is_animal,
                }
                for o in objects
            ])
            animals_json = json.dumps([
                {
                    "class": a.class_name,
                    "confidence": a.confidence,
                    "bbox": a.bbox,
                    "class_id": a.class_id,
                }
                for a in animals
            ])
            await asyncio.sleep(0.05)

            print(f"[INDEXING]   OCR...")
            ocr_results = run_ocr(full) if ext in SUPPORTED_IMAGE_TYPES else []
            ocr_json = json.dumps([
                {"text": r.text, "confidence": r.confidence, "bbox": r.bbox}
                for r in ocr_results
            ])
            await asyncio.sleep(0.05)

            # Face detection will run AFTER indexing completes in batch mode
            # Skip per-file face detection during indexing to improve performance
            face_instances = []
            faces_json = json.dumps([])

            print(f"[INDEXING]   Text extraction and embedding...")
            # Extract text from document files
            extracted_text = ""
            if ext in SUPPORTED_TEXT_TYPES:
                print(f"[INDEXING]   Extracting text from document...")
                extracted_text = extract_text_content(full)
                if extracted_text:
                    print(f"[INDEXING]   Extracted {len(extracted_text)} characters from document")
            
            # Build text corpus for embedding
            text_corpus = " ".join(
                [
                    os.path.basename(full),
                    " ".join(o.class_name for o in objects),
                    " ".join(r.text for r in ocr_results),
                    extracted_text,  # Add extracted text from documents
                ]
            )
            text_embed = get_sbert_embedding(text_corpus) or []
            await asyncio.sleep(0.05)

            # Convert AIF to WAV if needed
            file_to_store = full
            if ext in ['.aif', '.aiff']:
                print(f"[INDEXING]   Converting AIF to WAV...")
                file_to_store = convert_aif_to_wav(full)
                if file_to_store != full:
                    ext = '.wav'
                    print(f"[INDEXING]   ✓ Converted to WAV: {os.path.basename(file_to_store)}")

            record = {
                "path": os.path.abspath(file_to_store),
                "hash": file_hash,
                "type": ext,
                "date_taken": meta.get("date_taken"),
                "location": None,
                "size": stat.st_size,
                "width": meta.get("width"),
                "height": meta.get("height"),
                "camera": meta.get("camera"),
                "lens": meta.get("lens"),
                "text_embedding": to_blob(text_embed) if text_embed else None,
                "clip_embedding": to_blob(clip_embed) if clip_embed else None,
                "objects": objects_json,
                "faces": faces_json,
                "animals": animals_json,
                "text_content": text_corpus,
                "created_at": int(time.time()),
                "updated_at": int(time.time()),
            }
            
            # Store with fresh DB connection
            with get_db() as conn:
                media_id = upsert_media(conn, record)

                store_object_detections(
                    conn,
                    media_id,
                    [
                        {
                            "class_name": o.class_name,
                            "confidence": o.confidence,
                            "bbox": o.bbox,
                            "class_id": o.class_id,
                            "is_animal": o.is_animal,
                        }
                        for o in objects
                    ],
                )
                
                # Only store face embeddings for image files
                if ext in SUPPORTED_IMAGE_TYPES:
                    store_face_embeddings(
                        conn,
                        media_id,
                        [
                            {
                                "embedding": f["embedding"],
                                "bbox": f["bbox"],
                                "score": f["score"],
                            }
                            for f in face_instances
                        ],
                    )
                store_ocr_results(
                    conn,
                    media_id,
                    [
                        {"text": r.text, "confidence": r.confidence, "bbox": r.bbox}
                        for r in ocr_results
                    ],
                )

                # Store uncertain animal detections for review
                uncertain = [a for a in animals if a.confidence < 0.75]
                if uncertain:
                    store_uncertain_detections(conn, media_id, "animal", uncertain)

            if clip_embed:
                clip_embeddings.append(clip_embed)
                clip_ids.append(media_id)
            
            count += 1
            # Also increment aggregate counter across all sources
            indexing_state["processed_count"] = indexing_state.get("processed_count", 0) + 1
            
            total_faces += len(face_instances)
            total_animals += len(animals)
            indexing_state["processed"] = indexing_state.get("processed", 0) + 1
            indexing_state["faces_found"] = total_faces
            indexing_state["animals_found"] = total_animals
            
            # Update aggregate progress display
            total = indexing_state.get("total_files", total_files)
            processed = indexing_state.get("processed_count", 0)
            percentage = int((processed / max(1, total)) * 100)
            remaining = max(0, total - processed)
            
            # Update percentage in state for logs
            indexing_state["percentage"] = percentage
            
            # Show breakdown by type
            by_type = indexing_state.get("by_type", {})
            type_info = f"Images: {by_type.get('images', 0)} | Videos: {by_type.get('videos', 0)} | Audio: {by_type.get('audio', 0)} | Docs: {by_type.get('text', 0)}"
            
            indexing_state["current_file"] = (
                f"📊 {processed}/{total} ({percentage}%) | {type_info} | "
                f"Remaining: {remaining} | Processing: {name}"
            )
            
            print(f"[INDEXING] ✓ {name} ({processed}/{total}, {percentage}%)")
            
            # Clean up memory less frequently for speed
            if count % 10 == 0:
                gc.collect()
            # Minimal sleep for responsiveness
            await asyncio.sleep(0.01)

        except Exception as e:
            print(f"[INDEXING] ✗ Failed to index {name}: {e}")
            import traceback
            traceback.print_exc()
            indexing_state["processed"] = indexing_state.get("processed", 0) + 1
            indexing_state["processed_count"] = indexing_state.get("processed_count", 0) + 1
            if count % 10 == 0:
                gc.collect()
            await asyncio.sleep(0.01)

    if clip_embeddings:
        # Deduplicate: keep only first occurrence of each media_id
        seen = set()
        unique_embeddings = []
        unique_ids = []
        for emb, mid in zip(clip_embeddings, clip_ids):
            if mid not in seen:
                unique_embeddings.append(emb)
                unique_ids.append(mid)
                seen.add(mid)
        
        index = build_index(unique_embeddings, unique_ids, silo_name=current_silo_name)
        save_index(index, unique_ids, silo_name=current_silo_name)
    
    print(f"[INDEXING] ==========================================")
    print(f"[INDEXING] UNIFIED INDEXING COMPLETE")
    print(f"[INDEXING] Total processed: {count}")
    print(f"[INDEXING] Total already indexed: {already_indexed_count}")
    print(f"[INDEXING] Total newly indexed: {needs_indexing_count}")
    print(f"[INDEXING] Faces found: {total_faces}")
    print(f"[INDEXING] Animals found: {total_animals}")
    print(f"[INDEXING] ==========================================")
    
    return count

async def index_path(root: str, skip_videos: bool = False) -> int:
    """Backward compatibility wrapper - indexes a single path as a unified operation."""
    return await index_all_sources([root], skip_videos=skip_videos)


async def rebuild_faiss_index_from_db(silo_name: Optional[str] = None) -> int:
    """Rebuild FAISS index from existing embeddings in database.
    
    Args:
        silo_name: The silo to rebuild index for. If None, uses current silo context.
    """
    from .search_index import build_index, save_index
    
    # Use current silo context if not specified
    if not silo_name:
        try:
            from . import main
            silo_name = main._currently_processing_silo
        except:
            silo_name = None
    
    with get_db() as conn:
        cur = conn.execute(
            "SELECT id, clip_embedding FROM media_files WHERE clip_embedding IS NOT NULL ORDER BY id"
        )
        rows = cur.fetchall()
    
    if not rows:
        return 0
    
    clip_ids = []
    clip_embeddings = []
    
    for media_id, embedding_blob in rows:
        try:
            embedding = from_blob(embedding_blob)
            clip_embeddings.append(embedding)
            clip_ids.append(media_id)
        except Exception as e:
            print(f"Failed to load embedding for {media_id}: {e}")
    
    if clip_embeddings:
        index = build_index(clip_embeddings, clip_ids, silo_name=silo_name)
        save_index(index, clip_ids, silo_name=silo_name)
        print(f"[INDEXER] Rebuilt FAISS index for silo '{silo_name}' with {len(clip_ids)} embeddings")
        return len(clip_ids)
    return 0


def scan_all_sources(media_paths: list, skip_videos: bool = False) -> dict:
    """
    Scan all media sources upfront to get comprehensive statistics.
    Returns: {
        'total_files': int,
        'by_type': {'images': count, 'videos': count, 'audio': count, 'text': count},
        'files_by_path': {path: [list of files]},
        'already_indexed': int (from DB),
        'remaining': int
    }
    """
    from .db import get_db
    
    print(f"[PRE_SCAN] Scanning all {len(media_paths)} source(s) for comprehensive statistics...")
    
    all_files_by_path = {}
    file_counts = {"images": 0, "videos": 0, "audio": 0, "text": 0}
    total_files = 0
    
    # Scan all paths
    for media_path in media_paths:
        if not os.path.exists(media_path):
            print(f"[PRE_SCAN] ⚠ Path does not exist: {media_path}")
            continue
        
        path_files = []
        for dirpath, _, filenames in os.walk(media_path):
            for name in filenames:
                full = os.path.join(dirpath, name)
                ext = os.path.splitext(name)[1].lower()
                
                # Categorize by type
                if ext in SUPPORTED_IMAGE_TYPES:
                    file_counts["images"] += 1
                elif not skip_videos and ext in SUPPORTED_VIDEO_TYPES:
                    file_counts["videos"] += 1
                elif ext in SUPPORTED_AUDIO_TYPES:
                    file_counts["audio"] += 1
                elif ext in SUPPORTED_TEXT_TYPES:
                    file_counts["text"] += 1
                else:
                    continue  # Skip unsupported types
                
                path_files.append(full)
                total_files += 1
        
        if path_files:
            all_files_by_path[media_path] = path_files
            print(f"[PRE_SCAN] {media_path}: {len(path_files)} media files")
    
    # Count already indexed in database
    already_indexed = 0
    try:
        with get_db() as conn:
            cur = conn.execute("SELECT COUNT(*) as count FROM media_files")
            row = cur.fetchone()
            already_indexed = row[0] if row else 0
    except Exception as e:
        print(f"[PRE_SCAN] Could not count indexed files: {e}")
    
    result = {
        "total_files": total_files,
        "by_type": file_counts,
        "files_by_path": all_files_by_path,
        "already_indexed": already_indexed,
        "remaining": max(0, total_files - already_indexed)
    }
    
    print(f"[PRE_SCAN] ✓ Summary:")
    print(f"[PRE_SCAN]   Total files found: {total_files}")
    print(f"[PRE_SCAN]   - Images: {file_counts['images']}")
    print(f"[PRE_SCAN]   - Videos: {file_counts['videos']}")
    print(f"[PRE_SCAN]   - Audio: {file_counts['audio']}")
    print(f"[PRE_SCAN]   - Documents: {file_counts['text']}")
    print(f"[PRE_SCAN]   Already in DB: {already_indexed}")
    print(f"[PRE_SCAN]   Remaining to index: {result['remaining']}")
    
    return result


async def full_reindex() -> int:
    """Re-index all configured sources with progress tracking."""
    from .main import _get_silo_indexing_state, _currently_processing_silo
    from .face_cluster import detect_faces, cluster_faces, load_faces_from_db
    from .silo_manager import SiloManager
    
    indexing_state = _get_silo_indexing_state()
    cfg = load_config()
    
    # Get the silo name - CRITICAL: check if one is currently being processed
    current_silo_name = _currently_processing_silo
    if not current_silo_name:
        active_silo = SiloManager.get_active_silo()
        current_silo_name = active_silo.get("name", "default") if active_silo else "default"
    
    print(f"[INDEXING] ✓ Using silo: {current_silo_name}", flush=True)
    print(f"[INDEXING] ✓ Processing silo set: {_currently_processing_silo}", flush=True)
    
    media_paths = SiloManager.get_silo_media_paths(current_silo_name)
    
    # CRITICAL: Do NOT fall back to global config - each silo must have its own media_paths
    # Silos are isolated collections and should not share source folders
    if not media_paths:
        print(f"[INDEXING] No media paths configured for silo '{current_silo_name}'. User must add sources via UI.")
        indexing_state["status"] = "complete"
        indexing_state["message"] = "No media sources configured for this silo. Add a source folder to begin."
        return 0
    
    print(f"[INDEXING] Using media paths for silo '{current_silo_name}': {media_paths}")
    
    # PRE-SCAN all sources to get comprehensive statistics
    print(f"[INDEXING] ==========================================")
    print(f"[INDEXING] PHASE 1: SCANNING ALL SOURCES")
    print(f"[INDEXING] ==========================================")
    
    scan_stats = scan_all_sources(media_paths, skip_videos=cfg["processing"]["skip_videos"])
    
    # Store comprehensive stats in indexing state
    indexing_state["total_files"] = scan_stats["total_files"]
    indexing_state["by_type"] = scan_stats["by_type"]
    indexing_state["already_indexed"] = scan_stats["already_indexed"]
    indexing_state["remaining"] = scan_stats["remaining"]
    indexing_state["processed_count"] = 0
    
    # Display formatted stats
    stats_msg = (
        f"📊 Database: {scan_stats['already_indexed']} indexed | "
        f"Found: {scan_stats['total_files']} files | "
        f"Remaining: {scan_stats['remaining']} | "
        f"Images: {scan_stats['by_type']['images']} | "
        f"Videos: {scan_stats['by_type']['videos']} | "
        f"Audio: {scan_stats['by_type']['audio']} | "
        f"Documents: {scan_stats['by_type']['text']}"
    )
    indexing_state["current_file"] = stats_msg
    indexing_state["percentage"] = 0
    
    print(f"[INDEXING] {stats_msg}")
    print(f"[INDEXING] ==========================================")
    print(f"[INDEXING] PHASE 2: INDEXING FILES")
    print(f"[INDEXING] ==========================================")
    
    total_processed = 0
    
    # Index all sources in one unified operation maintaining aggregate progress
    try:
        print(f"[INDEXING] Starting unified indexing of all {len(media_paths)} sources")
        total_processed = await index_all_sources(media_paths, skip_videos=cfg["processing"]["skip_videos"])
        print(f"[INDEXING] Unified indexing complete: {total_processed} files processed")
    except Exception as e:
        print(f"[INDEXING] Error during unified indexing: {e}")
        import traceback
        traceback.print_exc()
        indexing_state["error"] = str(e)
    
    # Summary of media indexing phase
    print(f"[INDEXING] ==========================================")
    print(f"[INDEXING] ✅ MEDIA INDEXING COMPLETE")
    print(f"[INDEXING] Total files found: {scan_stats['total_files']}")
    print(f"[INDEXING] Breakdown:")
    print(f"[INDEXING]   - Images: {scan_stats['by_type']['images']}")
    print(f"[INDEXING]   - Videos: {scan_stats['by_type']['videos']}")
    print(f"[INDEXING]   - Audio: {scan_stats['by_type']['audio']}")
    print(f"[INDEXING]   - Documents: {scan_stats['by_type']['text']}")
    print(f"[INDEXING] Already indexed: {scan_stats['already_indexed']}")
    print(f"[INDEXING] Newly indexed: {total_processed}")
    print(f"[INDEXING] ==========================================")
    
    # Mark indexing as complete - face detection will be auto-triggered by frontend
    indexing_state["status"] = "complete"
    indexing_state["percentage"] = 100
    indexing_state["current_file"] = "✅ Media indexing complete! Face detection will start automatically..."
    
    print(f"[INDEXING] Media indexing complete - face detection will be auto-triggered")
    
    return total_processed


class WatchHandler(FileSystemEventHandler):
    def __init__(self, callback):
        super().__init__()
        self.callback = callback

    def on_modified(self, event):
        if not event.is_directory:
            asyncio.create_task(self.callback(event.src_path))

    def on_created(self, event):
        if not event.is_directory:
            asyncio.create_task(self.callback(event.src_path))


async def watch_directories(callback):
    # CRITICAL SECURITY: This function uses global config and would watch paths across ALL silos
    # File watching is DISABLED to prevent cross-silo data leakage
    # Each silo must be isolated - watching should be per-silo if implemented
    raise NotImplementedError(
        "File watching is disabled for security. Global config paths would leak data between silos. "
        "Implement per-silo watching if needed."
    )
    # REMOVED UNSAFE CODE:
    # cfg = load_config()
    # paths = cfg["storage"]["media_paths"]  # <-- GLOBAL PATHS! SECURITY BREACH!
    # observer = Observer()
    # handler = WatchHandler(callback)
    # for p in paths:
    #     observer.schedule(handler, p, recursive=True)
    #     observer.start()
    #     try:
    #         while True:
    #             await asyncio.sleep(1)
    #     finally:
    #         observer.stop()
    #         observer.join()


async def process_single(path: str):
    print(f"[PROCESS_SINGLE] Starting: {path}")
    
    # DEBUG: Check which silo context we have
    try:
        from . import main
        current_silo = main._currently_processing_silo
        print(f"[PROCESS_SINGLE] Current silo context: {current_silo}", flush=True)
    except:
        print(f"[PROCESS_SINGLE] Could not read silo context", flush=True)
    
    # DEBUG: Check what database path will be used
    try:
        from .db import get_db_path
        db_path = get_db_path()
        print(f"[PROCESS_SINGLE] Will use database: {db_path}", flush=True)
    except Exception as e:
        print(f"[PROCESS_SINGLE] ERROR getting db path: {e}", flush=True)
    
    cfg = load_config()
    if not is_media(path, cfg["processing"]["skip_videos"]):
        print(f"[PROCESS_SINGLE] Not a media file: {path}")
        return
    
    print(f"[PROCESS_SINGLE] File IS media, getting database connection...", flush=True)
    with get_db() as conn:
        print(f"[PROCESS_SINGLE] ✓ Got database connection successfully", flush=True)
        try:
            stat = os.stat(path)
            meta = extract_exif(path)
            ext = os.path.splitext(path)[1].lower()
            clip_embed = get_image_embedding(path) if ext in SUPPORTED_IMAGE_TYPES else None

            # Run object detection in thread pool to avoid blocking event loop
            if ext in SUPPORTED_IMAGE_TYPES:
                objects = await asyncio.to_thread(detect_objects, [path])
            else:
                objects = []
            animals = [o for o in objects if o.is_animal]
            objects_json = json.dumps([
                {
                    "class": o.class_name,
                    "confidence": o.confidence,
                    "bbox": o.bbox,
                    "class_id": o.class_id,
                    "is_animal": o.is_animal,
                }
                for o in objects
            ])
            animals_json = json.dumps([
                {
                    "class": a.class_name,
                    "confidence": a.confidence,
                    "bbox": a.bbox,
                    "class_id": a.class_id,
                }
                for a in animals
            ])

            ocr_results = run_ocr(path) if ext in SUPPORTED_IMAGE_TYPES else []
            ocr_json = json.dumps([
                {"text": r.text, "confidence": r.confidence, "bbox": r.bbox}
                for r in ocr_results
            ])

            # Skip face detection during indexing - run separately via /api/detect-faces-batch
            # face detection is slow and can cause hangs during large indexing operations
            face_instances = []
            faces_json = json.dumps([])

            text_corpus = " ".join(
                [
                    os.path.basename(path),
                    " ".join(o.class_name for o in objects),
                    " ".join(r.text for r in ocr_results),
                ]
            )
            text_embed = get_sbert_embedding(text_corpus) or []

            record = {
                "path": os.path.abspath(path),
                "hash": md5sum(path),
                "type": ext,
                "date_taken": meta.get("date_taken"),
                "location": None,
                "size": stat.st_size,
                "width": meta.get("width"),
                "height": meta.get("height"),
                "camera": meta.get("camera"),
                "lens": meta.get("lens"),
                "text_embedding": to_blob(text_embed) if text_embed else None,
                "clip_embedding": to_blob(clip_embed) if clip_embed else None,
                "objects": objects_json,
                "faces": faces_json,
                "animals": animals_json,
                "text_content": text_corpus,
                "created_at": int(time.time()),
                "updated_at": int(time.time()),
            }
            
            print(f"[PROCESS_SINGLE] About to call upsert_media with record for: {path}", flush=True)
            media_id = upsert_media(conn, record)
            print(f"[PROCESS_SINGLE] upsert_media returned media_id: {media_id}", flush=True)

            store_object_detections(
                conn,
                media_id,
                [
                    {
                        "class_name": o.class_name,
                        "confidence": o.confidence,
                        "bbox": o.bbox,
                        "class_id": o.class_id,
                        "is_animal": o.is_animal,
                    }
                    for o in objects
                ],
            )
            store_face_embeddings(
                conn,
                media_id,
                [
                    {
                        "embedding": f.embedding,
                        "bbox": f.bbox,
                        "score": f.score,
                    }
                    for f in face_instances
                ],
            )
            store_ocr_results(
                conn,
                media_id,
                [
                    {"text": r.text, "confidence": r.confidence, "bbox": r.bbox}
                    for r in ocr_results
                ],
            )

            uncertain = [a for a in animals if a.confidence < 0.75]
            if uncertain:
                store_uncertain_detections(conn, media_id, "animal", uncertain)

            conn.commit()
            print(f"[PROCESS_SINGLE_COMMITTED] Successfully committed: {path}")

            if clip_embed:
                # Get current silo context for index operations
                try:
                    from . import main
                    silo_name = main._currently_processing_silo
                except:
                    silo_name = None
                
                index, ids = load_index(len(clip_embed), silo_name=silo_name)
                vec = np.asarray([clip_embed], dtype="float32")
                faiss = __import__("faiss")
                faiss.normalize_L2(vec)
                index.add(vec)
                ids.append(media_id)
                save_index(index, ids, silo_name=silo_name)
        except Exception as e:
            print(f"[PROCESS_SINGLE_ERROR] DETAILED ERROR for {path}:")
            print(f"[PROCESS_SINGLE_ERROR] Exception type: {type(e).__name__}")
            print(f"[PROCESS_SINGLE_ERROR] Exception message: {str(e)}")
            import traceback
            traceback.print_exc()
            raise  # Re-raise so indexing loop knows there was an error


__all__ = ["full_reindex", "watch_directories", "process_single", "index_path", "is_media", "md5sum", "extract_exif", "SUPPORTED_IMAGE_TYPES", "SUPPORTED_AUDIO_TYPES"]
